"""
Presentación de avances — diseño limpio y natural.
conda run -n data_mining python3 generar_presentacion.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paleta minimalista
VERDE      = RGBColor(0x2E, 0x86, 0x48)   # verde oscuro, para acentos
VERDE_SUB  = RGBColor(0x5D, 0xAD, 0x72)   # verde medio, subtítulos
NEGRO      = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_TEXT  = RGBColor(0x44, 0x44, 0x44)
GRIS_SUAVE = RGBColor(0xF2, 0xF2, 0xF2)   # fondo de tabla alternado
GRIS_LINE  = RGBColor(0xCC, 0xCC, 0xCC)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
ROJO_TEXT  = RGBColor(0xB0, 0x2A, 0x2A)
AMARILLO   = RGBColor(0xF5, 0xA6, 0x23)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── utilidades ───────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0.5)):
    from pptx.util import Inches
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=GRIS_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def multiline(slide, lines, l, t, w, h, default_size=13,
              default_color=GRIS_TEXT, line_spacing=Pt(6)):
    """
    lines: list of (text, size, bold, color) or just str for defaults.
    """
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


def table(slide, headers, rows, l, t, w, h, font_size=12, col_widths=None):
    from pptx.util import Inches
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(h)
    ).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = VERDE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = BLANCO

    for ri, row in enumerate(rows):
        bg = GRIS_SUAVE if ri % 2 == 1 else BLANCO
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.color.rgb = NEGRO
    return tbl


def slide_title(slide, title, subtitle=None):
    """Línea verde fina + título grande."""
    rect(slide, 0, 0, 13.33, 0.07, fill=VERDE)
    txt(slide, title, 0.55, 0.25, 12.2, 0.8,
        size=28, bold=True, color=NEGRO)
    if subtitle:
        txt(slide, subtitle, 0.55, 0.95, 12.2, 0.4,
            size=13, color=VERDE_SUB, italic=True)


def divider(slide, y):
    rect(slide, 0.55, y, 12.2, 0.02, fill=GRIS_LINE)


def slide_num(slide, n, total=11):
    txt(slide, f"{n} / {total}", 12.6, 7.2, 0.6, 0.25,
        size=9, color=GRIS_LINE, align=PP_ALIGN.RIGHT)


# ── 1. Portada ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=BLANCO)
rect(s, 0, 0, 0.18, 7.5, fill=VERDE)          # barra lateral verde
rect(s, 0.18, 3.5, 13.15, 0.04, fill=GRIS_LINE)

txt(s, "Avances de Estancia", 0.65, 1.6, 12.0, 1.1,
    size=44, bold=True, color=NEGRO)
txt(s, "Semanas 1 y 2  —  Gold Standard para Análisis de Sentimiento",
    0.65, 2.65, 12.0, 0.55, size=18, color=VERDE_SUB)

txt(s, "Hugo Emiliano Vargas Briones", 0.65, 3.85, 10.0, 0.45,
    size=14, bold=True, color=NEGRO)
txt(s, "Asesor: Dr. José Lázaro Martínez Rodríguez", 0.65, 4.3, 10.0, 0.4,
    size=13, color=GRIS_TEXT)
txt(s, "Correcaminos hacia la Ciencia 2026  ·  UAT FIC", 0.65, 4.7, 10.0, 0.4,
    size=13, color=GRIS_TEXT)
txt(s, "13 de junio de 2026", 0.65, 5.1, 10.0, 0.4,
    size=13, color=GRIS_TEXT)


# ── 2. Contexto ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Contexto del proyecto")
divider(s, 1.45)
slide_num(s, 2)

txt(s, "La tesis evalúa un ensamble de tres LLMs locales que clasifica el sentimiento de reseñas "
       "de productos agrícolas de Amazon (semillas y pastos, en inglés) en positivo, neutral o negativo.",
    0.55, 1.6, 12.2, 0.9, size=14, color=GRIS_TEXT, wrap=True)

# tres columnas: modelos / etiquetado / gap
for i, (titulo, cuerpo) in enumerate([
    ("Modelos en el ensamble",
     "Qwen2.5-1.5b\nGemma2-2b\nPhi4-mini\n(Ollama, local)"),
    ("Etiquetado",
     "Votación mayoritaria\n→ positivo / neutral / negativo\n\n58,108 reseñas ya etiquetadas"),
    ("Lo que faltaba",
     "Una referencia humana\npara saber si esas etiquetas\nson confiables.\n\nEso es el gold standard."),
]):
    x = 0.55 + i * 4.2
    rect(s, x, 2.7, 3.9, 4.4, fill=GRIS_SUAVE)
    rect(s, x, 2.7, 3.9, 0.45, fill=VERDE)
    txt(s, titulo, x + 0.15, 2.72, 3.6, 0.4,
        size=12, bold=True, color=BLANCO)
    txt(s, cuerpo, x + 0.15, 3.25, 3.6, 3.7,
        size=13, color=NEGRO if i < 2 else ROJO_TEXT, wrap=True)


# ── 3. Distribución del dataset ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Distribución del dataset",
            "58,108 reseñas — primer análisis antes de diseñar el muestreo")
divider(s, 1.45)
slide_num(s, 3)

table(s,
    ["Estadístico", "Palabras"],
    [["Mediana", "20"], ["Media", "35.1"],
     ["10% tiene ≤", "4 palabras"], ["90% tiene ≤", "82 palabras"], ["Máxima", "338"]],
    0.55, 1.6, 3.7, 2.9,
    col_widths=[2.2, 1.4], font_size=12)

table(s,
    ["Clase", "Reseñas", "%"],
    [["Positivo", "26,614", "45.8%"],
     ["Negativo", "25,558", "44.0%"],
     ["Neutral",  " 5,936", "10.2%"]],
    4.7, 1.6, 4.3, 2.1,
    col_widths=[1.8, 1.4, 1.0], font_size=12)

table(s,
    ["Concordancia LLMs", "Reseñas", "%"],
    [["1.0 — todos coinciden",  "45,692", "78.6%"],
     ["0.67 — 2 de 3 coinciden","11,164", "19.2%"],
     ["0.33 — ninguno coincide"," 1,252",  "2.2%"]],
    4.7, 3.85, 8.05, 2.1,
    col_widths=[4.2, 1.8, 1.8], font_size=12)

multiline(s, [
    ("Distribución sesgada a la derecha.", 12, True, NEGRO),
    "Mediana 20 palabras, pero media de 35. Las reseñas cortas son mayoría.",
    "Y son las que más le cuestan al pipeline.",
], 0.55, 4.65, 3.7, 1.8)

txt(s, "El 78.6% de los casos son fáciles (todos coinciden).\n"
       "Un gold standard construido solo con esos casos\ntendría sesgo de selección.",
    4.7, 6.05, 7.9, 1.2, size=12, color=GRIS_TEXT, wrap=True)


# ── 4. Hallazgo ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Hallazgo: las reseñas neutras son más cortas y más ambiguas",
            "Cruce longitud × clase × concordancia")
divider(s, 1.45)
slide_num(s, 4)

table(s,
    ["Clase", "Mediana de palabras"],
    [["Positivo", "21"], ["Negativo", "21"], ["Neutral", "16  (−24%)"]],
    0.55, 1.6, 4.8, 2.0,
    col_widths=[2.4, 2.3], font_size=13)

table(s,
    ["Concordancia", "% con ≤5 palabras"],
    [["conc=1.0 (todos coinciden)", "13.4%"],
     ["conc=0.33 (nadie coincide)", "25.9%  (×2)"]],
    0.55, 3.85, 4.8, 1.7,
    col_widths=[2.8, 1.9], font_size=13)

txt(s, "En los casos más difíciles hay el doble de reseñas ultracortas.",
    0.55, 5.65, 4.8, 0.6, size=12, italic=True, color=GRIS_TEXT)

# bloque derecho
txt(s, "Por qué falla el pipeline en neutral", 6.0, 1.6, 6.8, 0.45,
    size=14, bold=True, color=NEGRO)
txt(s, "Las reseñas cortas tienen sentimiento implícito. La votación mayoritaria "
       "entre tres modelos no alcanza a desambiguarlas.",
    6.0, 2.1, 6.8, 1.0, size=13, color=GRIS_TEXT, wrap=True)

table(s,
    ["Reseña (inglés)", "LLMs dicen", "Humano dice"],
    [["\"cat loves it\"",    "Neutral", "Positivo"],
     ["\"No complaints.\"",  "Neutral", "Positivo"],
     ["\"My dog liked it\"", "Neutral", "Positivo"]],
    6.0, 3.3, 6.8, 2.1,
    col_widths=[3.2, 1.7, 1.8], font_size=13)

txt(s, "Los tres son reseñas que cualquier humano lee como positivas.\nEl ensamble no lo ve.",
    6.0, 5.55, 6.8, 0.9, size=12, italic=True, color=ROJO_TEXT, wrap=True)


# ── 5. Pool disponible ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Restricción del pool — clave para el diseño del muestreo",
            "Cuántas reseñas hay disponibles por clase × concordancia")
divider(s, 1.45)
slide_num(s, 5)

table(s,
    ["Clase", "conc = 1.0", "conc = 0.67", "conc = 0.33"],
    [["Positivo", "22,457", "4,157",  "0"],
     ["Neutral",     "610", "4,074", "1,252"],
     ["Negativo", "22,625", "2,933",  "0"]],
    1.4, 1.65, 10.4, 2.4,
    col_widths=[2.5, 2.5, 2.6, 2.6], font_size=15)

divider(s, 4.35)

txt(s, "Por que positivo y negativo tienen cero casos con conc=0.33",
    0.55, 4.5, 12.2, 0.45, size=14, bold=True, color=NEGRO)

txt(s, "Cuando los tres modelos discrepan totalmente, el ensamble no puede resolver mayoría "
       "y asigna Neutral por defecto. Los 1,252 casos de conc=0.33 son todos Neutral — "
       "no por el contenido de las reseñas, sino por la mecánica del sistema de votación.",
    0.55, 5.05, 12.2, 1.5, size=13, color=GRIS_TEXT, wrap=True)


# ── 6. Diseño del gold standard ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Diseño del gold standard — 999 reseñas estratificadas",
            "Semana 1")
divider(s, 1.45)
slide_num(s, 6)

txt(s, "Un gold standard con solo casos de conc=1.0 (los obvios) siempre va a hacer "
       "que el pipeline se vea bien. Los errores reales aparecen en los casos donde "
       "los modelos discrepan.",
    0.55, 1.6, 12.2, 1.0, size=13, color=GRIS_TEXT, wrap=True)

table(s,
    ["Concordancia", "Proporción", "Reseñas / clase", "Justificación"],
    [["1.0  — fáciles",       "40%", "~133", "Casos ancla, calibran al anotador"],
     ["0.67 — intermedios",   "50%", "~167", "Zona de ambigüedad real"],
     ["0.33 — difíciles",     "10%", "~34 (solo Neutral)", "Máxima discrepancia entre modelos"]],
    0.55, 2.75, 12.2, 2.6,
    col_widths=[2.8, 1.8, 2.6, 4.8], font_size=13)

divider(s, 5.55)
txt(s, "Resultado:  999 reseñas  |  333 por clase  |  archivo: ANOTAR/para_anotar_humano_es.csv",
    0.55, 5.7, 12.2, 0.5, size=14, bold=True, color=VERDE)
txt(s, "Cada reseña tiene texto en español (traducción) y texto original en inglés.",
    0.55, 6.2, 12.2, 0.4, size=12, color=GRIS_TEXT, italic=True)


# ── 7. Protocolo + entregables Semana 1 ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Protocolo de anotación y entregables de Semana 1")
divider(s, 1.45)
slide_num(s, 7)

# izq
txt(s, "Regla principal del protocolo", 0.55, 1.6, 5.8, 0.4,
    size=13, bold=True, color=NEGRO)
multiline(s, [
    "El texto de la reseña manda sobre el rating numérico.",
    "",
    "La estrella (1–5) es contexto, no criterio de clasificación.",
    "",
    "Empate entre anotadores  →  Neutral por defecto.",
    "",
    "Ejemplo: reseña de 5 estrellas con queja explícita  →  Negativo.",
], 0.55, 2.1, 5.8, 4.8, default_size=13)

divider(s, 1.6)   # vertical line simulation not needed — just space

# der
txt(s, "Archivos entregados", 7.0, 1.6, 5.8, 0.4,
    size=13, bold=True, color=NEGRO)
multiline(s, [
    ("protocolo_anotacion.html", 13, True, NEGRO),
    "Reglas completas, imprimible.",
    "",
    ("para_anotar_humano.csv", 13, True, NEGRO),
    "Solo índice + texto. Sin etiquetas LLM visibles (evita sesgo).",
    "",
    ("para_anotar_humano_es.csv", 13, True, NEGRO),
    "999 reseñas en español (traducción vía Claude Sonnet 4.6).",
    "",
    ("crear_formulario_v2_parte01…10.gs", 13, True, NEGRO),
    "10 scripts para Google Forms, listos para ejecutar.",
], 7.0, 2.1, 6.1, 5.2, default_size=13)

rect(s, 6.6, 1.45, 0.03, 5.8, fill=GRIS_LINE)  # separador vertical


# ── 8. Herramienta de anotación ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Herramienta de anotación — 10 Google Forms",
            "Semana 2")
divider(s, 1.45)
slide_num(s, 8)

txt(s, "El problema", 0.55, 1.6, 12.2, 0.4, size=13, bold=True, color=NEGRO)
txt(s, "Un formulario con 999 preguntas en una sola página deja de responder en el navegador, "
       "incluso con hardware M4. La solución fue dividir en 10 formularios de ~100 reseñas cada uno.",
    0.55, 2.05, 12.2, 0.9, size=13, color=GRIS_TEXT, wrap=True)

divider(s, 3.15)

txt(s, "Estructura de cada formulario", 0.55, 3.3, 6.0, 0.4,
    size=13, bold=True, color=NEGRO)
multiline(s, [
    "Una página por reseña (sin scroll infinito).",
    "",
    "Cada página muestra:",
    "  · Reseña #N  ·  ID original",
    "  · Texto en español",
    "  · Texto en inglés (original)",
    "  · Selección: Positivo / Neutral / Negativo  (obligatorio)",
    "  · Justificación (opcional)",
], 0.55, 3.8, 5.9, 3.5, default_size=13)

txt(s, "Cómo se crea cada formulario", 7.2, 3.3, 5.7, 0.4,
    size=13, bold=True, color=NEGRO)
multiline(s, [
    "1. Abrir script.google.com",
    "2. Pegar contenido de crear_formulario_v2_parteXX.gs",
    "3. Ejecutar crearFormulario()",
    "4. Aceptar permisos → el formulario aparece en Google Drive",
], 7.2, 3.8, 5.7, 3.0, default_size=13)


# ── 9. Primeras respuestas ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Primeras respuestas — Parte 1 completada",
            "Semana 2  ·  100 reseñas  ·  3 anotadores")
divider(s, 1.45)
slide_num(s, 9)

table(s,
    ["Parte", "Reseñas", "Estado", "Anotadores"],
    [["Parte 1  (reseñas 1–100)", "100", "Completada", "3 respuestas"],
     ["Partes 2–10",              "899", "Pendiente",  "—"]],
    0.55, 1.6, 12.2, 1.9,
    col_widths=[5.0, 1.8, 2.4, 2.8], font_size=13)

txt(s, "Variedad de acuerdo en las primeras 100 reseñas", 0.55, 3.7, 12.2, 0.45,
    size=13, bold=True, color=NEGRO)

table(s,
    ["Reseña", "Resultado (3 anotadores)", "Tipo de caso", "Etiqueta final"],
    [["#58  ID 46152",
      "100%  Positivo",
      "Acuerdo total",
      "Positivo"],
     ["#52  ID 4625",
      "66.7% Positivo  /  33.3% Negativo",
      "Mayoría clara",
      "Positivo"],
     ["#45  ID 17134",
      "33.3% cada clase",
      "Empate",
      "Neutral (protocolo)"]],
    0.55, 4.2, 12.2, 2.4,
    col_widths=[2.4, 4.4, 2.6, 2.6], font_size=12)

txt(s, "No es coincidencia — el muestreo fue diseñado para cubrir los tres niveles.",
    0.55, 6.75, 12.2, 0.4,
    size=12, italic=True, color=VERDE_SUB)


# ── 10. Próximos pasos ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Próximos pasos — Semana 3  (15–20 jun)")
divider(s, 1.45)
slide_num(s, 10)

pasos = [
    ("1", "Completar la anotación humana",
     "Partes 2–10  →  899 reseñas restantes"),
    ("2", "Calcular métricas de evaluación",
     "Cohen's Kappa  ·  F1-Macro por clase  ·  % concordancia por modelo individual"),
    ("3", "Análisis de errores",
     "Clasificar los fallos por tipo: longitud, sentimiento implícito, ambigüedad"),
    ("4", "Ajuste del esquema de etiquetado v2.0",
     "Si los datos revelan clases conflictivas, refinar el protocolo de anotación"),
]

for i, (num, titulo, desc) in enumerate(pasos):
    y = 1.65 + i * 1.3
    rect(s, 0.55, y, 0.55, 0.55, fill=VERDE)
    txt(s, num, 0.55, y + 0.03, 0.55, 0.5,
        size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(s, titulo, 1.3, y, 11.4, 0.38,
        size=14, bold=True, color=NEGRO)
    txt(s, desc, 1.3, y + 0.38, 11.4, 0.4,
        size=12, color=GRIS_TEXT)


# ── 11. Resumen de entregables ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_title(s, "Resumen de entregables", "Estado al 13 de junio de 2026")
divider(s, 1.45)
slide_num(s, 11)

filas = [
    ("Gold standard diseñado — 999 reseñas estratificadas", "Semana 1", True),
    ("Protocolo de anotación v1.0",                          "Semana 1", True),
    ("999 reseñas traducidas al español",                    "Semana 1", True),
    ("10 Google Forms listos para anotación",                "Semana 2", True),
    ("Parte 1 anotada  (100 reseñas, 3 anotadores)",         "Semana 2", True),
    ("Partes 2–10 anotadas  (899 reseñas)",                  "Semana 2–3", False),
    ("Métricas: kappa, F1-macro, análisis de errores",       "Semana 3", False),
    ("Reporte final y video",                                 "Semana 4", False),
]

n = len(filas) + 1
tbl = s.shapes.add_table(
    n, 3, Inches(0.55), Inches(1.65), Inches(12.2), Inches(5.6)
).table
tbl.columns[0].width = Inches(8.2)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.0)

for ci, hdr in enumerate(["Entregable", "Semana", "Estado"]):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NEGRO
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = hdr
    run.font.name = "Calibri"
    run.font.size = Pt(13); run.font.bold = True
    run.font.color.rgb = BLANCO

for ri, (entregable, semana, done) in enumerate(filas):
    bg = GRIS_SUAVE if ri % 2 == 1 else BLANCO
    estado_txt = "Completado" if done else "Pendiente"
    estado_col = VERDE if done else AMARILLO
    for ci, val in enumerate([entregable, semana, estado_txt]):
        cell = tbl.cell(ri + 1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = val
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = (ci == 2)
        run.font.color.rgb = estado_col if ci == 2 else NEGRO


# ── guardar ───────────────────────────────────────────────────────────────────
OUT = "/Users/emiliano/Desktop/Proyecto de Tesis/estancia/Avances_Estancia_Emiliano.pptx"
prs.save(OUT)
print(f"Guardado: {OUT}")
print(f"Diapositivas: {len(prs.slides)}")
